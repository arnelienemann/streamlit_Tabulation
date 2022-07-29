from configparser import Interpolation
from ctypes import util
from pickle import TRUE
from re import X
from turtle import color
from h11 import Data
from matplotlib import image
from pyparsing import col
import streamlit as st
import numpy as np
import pandas as pd
from wordcloud import STOPWORDS, WordCloud
import matplotlib.pyplot as plt
import datetime
import io
from PIL import Image 
import PIL 

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import ColorFormat, RGBColor

st.header("Reporting Automation")

file_path = "data-reportautomation/style.pptx"
prs = Presentation(file_path)
title_slide_layout = prs.slide_layouts[0]
divider_slide_layout = prs.slide_layouts[3]
content_slide_layout = prs.slide_layouts[2]
text_slide_layout = prs.slide_layouts[5]
end_slide_layout = prs.slide_layouts[19]

data_file = st.file_uploader("Upload .csv file:")

if data_file is not None:

    df_data = pd.read_csv(data_file, sep=";")

    st.write("Project Data:")
    st.dataframe(df_data)


    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = "Survey Results"
    title_slide.shapes[1].text = "Arne Lienemann \nConsumer Insights | Global Portfolio Management"

    divider_slide = prs.slides.add_slide(divider_slide_layout)
    divider_slide.shapes.title.text = "Results"

    for column in df_data.columns:
        
        df_data = df_data[df_data != 999]
                    
        absolute_frequencies = df_data[column].value_counts(normalize=False)
        sum = absolute_frequencies.sum()

        relative_frequencies = absolute_frequencies/sum
        relative_frequencies.sort_index(inplace=True)

        
        st.write(column + " (n = " + str(sum) + ")")
        st.dataframe(relative_frequencies)
        #relative_frequencies = pd.DataFrame(relative_frequencies)

        content_slide = prs.slides.add_slide(content_slide_layout)
        
        #Find max value to print key take away
        #key_take_away = relative_frequencies.idxmax()[0]
        #content_slide.shapes.title.text = "Results" + "\n" + key_take_away

        content_slide.shapes[1].text = column + " (n = " + str(sum) + ")."
        #content_slide.shapes[2].text = "" #str(index + 3)
        #content_slide.shapes[3].text = Project_name

        chart_data = CategoryChartData()
        chart_data.categories = relative_frequencies.index
        chart_data.add_series("", relative_frequencies)

        graphic_frame = content_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.73), Inches(1.35), Inches(9.89), Inches(5.39), chart_data)
        graphic_frame.chart.series[0].format.fill.solid()
        graphic_frame.chart.series[0].format.fill.fore_color.rgb = RGBColor(153, 174, 181)
        graphic_frame.chart.value_axis.visible = False
        graphic_frame.chart.category_axis.tick_labels.font.size = Pt(10.5)
        graphic_frame.chart.plots[0].has_data_labels = True
        graphic_frame.chart.plots[0].data_labels.font.size = Pt(10.5)
        graphic_frame.chart.plots[0].data_labels.number_format = '0%'
    
    end_slide = prs.slides.add_slide(end_slide_layout)
    end_slide.shapes.title.text = "Thank you"
    end_slide.shapes[1].text ="Contact \nArne Lienemann \nConsumer Insights | Global Portfolio Management \nArne.Lienemann@Sennheiser-ce.com"

    prs.save('data-reportautomation/newppt.pptx')

    #with open(file_path, 'rb') as my_file:
    with open('data-reportautomation/newppt.pptx', 'rb') as my_file:
        st.download_button(label = 'Download', data = my_file, file_name = 'results.pptx') 


    #python -m streamlit run reportautomation.py